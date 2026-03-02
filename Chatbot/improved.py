import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import threading
import time
import datetime
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from langchain.vectorstores import FAISS as LangchainFAISS
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.llms import Ollama
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter

class ModernDocumentChatGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("🤖 AI Document Chat Assistant")
        
        # Get screen dimensions and calculate optimal window size
        screen_width = self.root.winfo_screenwidth()
        screen_height = self.root.winfo_screenheight()
        
        # Calculate window size based on screen resolution
        if screen_width >= 1920:  # High resolution monitors
            window_width = min(1400, int(screen_width * 0.7))
            window_height = min(1000, int(screen_height * 0.8))
        elif screen_width >= 1366:  # Standard HD monitors
            window_width = min(1200, int(screen_width * 0.8))
            window_height = min(900, int(screen_height * 0.85))
        else:  # Smaller screens
            window_width = min(1000, int(screen_width * 0.9))
            window_height = min(700, int(screen_height * 0.9))
        
        self.root.geometry(f"{window_width}x{window_height}")
        self.root.configure(bg='#1a1a2e')
        
        # Make window resizable with minimum size
        self.root.minsize(800, 600)
        self.root.resizable(True, True)
        
        # Modern color scheme
        self.colors = {
            'bg_primary': '#1a1a2e',
            'bg_secondary': '#16213e',
            'bg_accent': '#0f3460',
            'text_primary': '#ffffff',
            'text_secondary': '#b8c6db',
            'accent_blue': '#00d4ff',
            'accent_green': '#00ff88',
            'accent_orange': '#ff6b35',
            'accent_purple': '#a663cc',
            'button_hover': '#533483',
            'error': '#ff4757',
            'success': '#2ed573'
        }
        
        # Initialize variables
        self.qa_system = None
        self.documents_loaded = False
        self.embed_model = None
        self.supported_files_count = 0
        self.response_times = []  # Track response times for analytics
        
        # Configure styles
        self.setup_styles()
        
        # Create GUI elements
        self.create_widgets()
        
    def setup_styles(self):
        """Configure modern ttk styles"""
        style = ttk.Style()
        style.theme_use('clam')
        
        # Configure main frame style
        style.configure('Main.TFrame', background=self.colors['bg_primary'])
        
        # Configure label frame styles
        style.configure('Card.TLabelFrame', 
                       background=self.colors['bg_secondary'],
                       foreground=self.colors['text_primary'],
                       borderwidth=2,
                       relief='flat')
        style.configure('Card.TLabelFrame.Label',
                       background=self.colors['bg_secondary'],
                       foreground=self.colors['accent_blue'],
                       font=('Segoe UI', 11, 'bold'))
        
        # Configure button styles
        style.configure('Accent.TButton',
                       background=self.colors['accent_blue'],
                       foreground='white',
                       font=('Segoe UI', 10, 'bold'),
                       borderwidth=0,
                       focuscolor='none')
        style.map('Accent.TButton',
                 background=[('active', self.colors['button_hover'])])
        
        style.configure('Success.TButton',
                       background=self.colors['accent_green'],
                       foreground='white',
                       font=('Segoe UI', 10, 'bold'),
                       borderwidth=0,
                       focuscolor='none')
        style.map('Success.TButton',
                 background=[('active', '#27ae60')])
        
        style.configure('Warning.TButton',
                       background=self.colors['accent_orange'],
                       foreground='white',
                       font=('Segoe UI', 10, 'bold'),
                       borderwidth=0,
                       focuscolor='none')
        style.map('Warning.TButton',
                 background=[('active', '#e55039')])
        
        # Configure label styles
        style.configure('Title.TLabel',
                       background=self.colors['bg_primary'],
                       foreground=self.colors['text_primary'],
                       font=('Segoe UI', 20, 'bold'))
        
        style.configure('Modern.TLabel',
                       background=self.colors['bg_secondary'],
                       foreground=self.colors['text_secondary'],
                       font=('Segoe UI', 10))
        
        style.configure('Status.TLabel',
                       background=self.colors['bg_secondary'],
                       foreground=self.colors['accent_green'],
                       font=('Segoe UI', 10, 'bold'))
        
        # Configure entry styles
        style.configure('Modern.TEntry',
                       fieldbackground=self.colors['bg_accent'],
                       foreground=self.colors['text_primary'],
                       borderwidth=2,
                       insertcolor=self.colors['text_primary'])
        
        # Configure progressbar style
        style.configure('Modern.Horizontal.TProgressbar',
                       background=self.colors['accent_blue'],
                       troughcolor=self.colors['bg_accent'],
                       borderwidth=0,
                       lightcolor=self.colors['accent_blue'],
                       darkcolor=self.colors['accent_blue'])
        
    def create_widgets(self):
        # Main frame with gradient effect
        main_frame = ttk.Frame(self.root, style='Main.TFrame', padding="20")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Configure grid weights
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.rowconfigure(4, weight=1)
        
        # Modern title with emoji
        title_frame = tk.Frame(main_frame, bg=self.colors['bg_primary'])
        title_frame.grid(row=0, column=0, columnspan=3, pady=(0, 30), sticky=(tk.W, tk.E))
        
        title_label = ttk.Label(title_frame, text="🤖 AI Document Chat Assistant", 
                               style='Title.TLabel')
        title_label.pack()
        
        subtitle_label = ttk.Label(title_frame, 
                                  text="Chat with your PDF, DOCX, PPTX, and XML documents using AI",
                                  background=self.colors['bg_primary'],
                                  foreground=self.colors['text_secondary'],
                                  font=('Segoe UI', 11))
        subtitle_label.pack(pady=(5, 0))
        
        # Document loading section with modern card design
        doc_frame = ttk.LabelFrame(main_frame, text="📁 Document Management")
        doc_frame.grid(row=1, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 20))
        doc_frame.columnconfigure(1, weight=1)
        
        # Folder selection with modern styling
        ttk.Label(doc_frame, text="📂 Documents Folder:", style='Modern.TLabel').grid(
            row=0, column=0, sticky=tk.W, padx=(0, 15), pady=(0, 10))
        
        self.folder_var = tk.StringVar()
        self.folder_entry = ttk.Entry(doc_frame, textvariable=self.folder_var, 
                                     state='readonly', style='Modern.TEntry', width=50)
        self.folder_entry.grid(row=0, column=1, sticky=(tk.W, tk.E), padx=(0, 15), pady=(0, 10))
        
        self.browse_btn = ttk.Button(doc_frame, text="🔍 Browse", 
                                    command=self.browse_folder, style='Accent.TButton')
        self.browse_btn.grid(row=0, column=2, pady=(0, 10))
        
        # Supported formats info
        formats_label = ttk.Label(doc_frame, 
                                 text="📋 Supported formats: PDF, DOCX, PPTX, XML",
                                 style='Modern.TLabel')
        formats_label.grid(row=1, column=0, columnspan=3, sticky=tk.W, pady=(0, 10))
        
        # Load documents button
        self.load_btn = ttk.Button(doc_frame, text="⚡ Load Documents", 
                                  command=self.load_documents_thread, 
                                  state='disabled', style='Success.TButton')
        self.load_btn.grid(row=2, column=0, columnspan=3, pady=(10, 0))
        
        # Modern status section
        status_frame = ttk.LabelFrame(main_frame, text="📊 System Status")
        status_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 20))
        status_frame.columnconfigure(0, weight=1)
        
        self.status_var = tk.StringVar(value="🟡 Ready - Please select a documents folder")
        self.status_label = ttk.Label(status_frame, textvariable=self.status_var, 
                                     style='Status.TLabel')
        self.status_label.grid(row=0, column=0, sticky=tk.W, pady=(0, 10))
        
        # Modern progress bar
        self.progress = ttk.Progressbar(status_frame, mode='indeterminate', 
                                       style='Modern.Horizontal.TProgressbar')
        self.progress.grid(row=1, column=0, sticky=(tk.W, tk.E))
        
        # Document count display
        self.doc_count_var = tk.StringVar(value="")
        self.doc_count_label = ttk.Label(status_frame, textvariable=self.doc_count_var,
                                        style='Modern.TLabel')
        self.doc_count_label.grid(row=2, column=0, sticky=tk.W, pady=(5, 0))
        
        # Performance metrics display
        self.perf_var = tk.StringVar(value="")
        self.perf_label = ttk.Label(status_frame, textvariable=self.perf_var,
                                   style='Modern.TLabel')
        self.perf_label.grid(row=3, column=0, sticky=tk.W, pady=(2, 0))
        
        # Modern chat section
        chat_frame = ttk.LabelFrame(main_frame, text="💬 AI Chat Interface")
        chat_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=(0, 20))
        chat_frame.columnconfigure(0, weight=1)
        chat_frame.rowconfigure(0, weight=1)
        
        # Chat display with modern styling and responsive height
        chat_container = tk.Frame(chat_frame, bg=self.colors['bg_accent'], bd=2, relief='flat')
        chat_container.grid(row=0, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=(0, 15))
        chat_container.columnconfigure(0, weight=1)
        chat_container.rowconfigure(0, weight=1)
        
        # Calculate chat height based on window size
        self.root.update_idletasks()
        window_height = self.root.winfo_height()
        chat_height = max(15, min(25, int(window_height / 40)))
        
        self.chat_display = scrolledtext.ScrolledText(
            chat_container, height=chat_height, state='disabled', wrap=tk.WORD, 
            font=('Consolas', 10), bg=self.colors['bg_accent'], 
            fg=self.colors['text_primary'], insertbackground=self.colors['text_primary'],
            selectbackground=self.colors['accent_blue'], selectforeground='white',
            bd=0, highlightthickness=0)
        self.chat_display.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S), padx=10, pady=10)
        
        # Query input section
        input_frame = tk.Frame(chat_frame, bg=self.colors['bg_secondary'])
        input_frame.grid(row=1, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        input_frame.columnconfigure(1, weight=1)
        
        ttk.Label(input_frame, text="❓ Ask a question:", style='Modern.TLabel').grid(
            row=0, column=0, sticky=tk.W, padx=(0, 15))
        
        self.query_var = tk.StringVar()
        self.query_entry = ttk.Entry(input_frame, textvariable=self.query_var, 
                                    state='disabled', style='Modern.TEntry', font=('Segoe UI', 11))
        self.query_entry.grid(row=0, column=1, sticky=(tk.W, tk.E), padx=(0, 15))
        self.query_entry.bind('<Return>', self.ask_question_event)
        
        self.ask_btn = ttk.Button(input_frame, text="🚀 Ask AI", 
                                 command=self.ask_question_thread, 
                                 state='disabled', style='Accent.TButton')
        self.ask_btn.grid(row=0, column=2)
        
        # Control buttons with performance stats
        control_frame = tk.Frame(chat_frame, bg=self.colors['bg_secondary'])
        control_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E))
        
        # Performance stats button
        self.stats_btn = ttk.Button(control_frame, text="📊 Performance Stats", 
                                   command=self.show_performance_stats, style='Accent.TButton')
        self.stats_btn.pack(side=tk.LEFT)
        
        self.clear_btn = ttk.Button(control_frame, text="🧹 Clear Chat", 
                                   command=self.clear_chat, style='Warning.TButton')
        self.clear_btn.pack(side=tk.RIGHT, padx=(10, 0))
        
        # Add welcome message
        self.add_welcome_message()
        
    def add_welcome_message(self):
        """Add a colorful welcome message to the chat"""
        welcome_msg = """🎉 Welcome to AI Document Chat Assistant!

✨ Features:
• 📄 PDF document support
• 📝 Word document support (.docx)
• 📊 PowerPoint presentation support (.pptx)
• 🗂️ XML file support
• 🤖 Powered by Llama 3.2 AI model

🚀 Getting Started:
1. Click 'Browse' to select your documents folder
2. Click 'Load Documents' to process your files
3. Start asking questions about your documents!

Ready to explore your documents with AI? Let's get started! 🌟"""
        
        self.add_to_chat("🤖 AI Assistant", welcome_msg, color='system')
        
    def browse_folder(self):
        folder = filedialog.askdirectory(title="Select Documents Folder")
        if folder:
            self.folder_var.set(folder)
            self.load_btn.config(state='normal')
            self.status_var.set(f"🟢 Folder selected: {os.path.basename(folder)}")
            
    def load_documents_thread(self):
        thread = threading.Thread(target=self.load_documents)
        thread.daemon = True
        thread.start()
        
    def load_documents(self):
        try:
            folder = self.folder_var.get()
            if not folder or not os.path.exists(folder):
                messagebox.showerror("Error", "Please select a valid documents folder")
                return
                
            # Update UI
            self.root.after(0, lambda: self.progress.start())
            self.root.after(0, lambda: self.load_btn.config(state='disabled'))
            self.root.after(0, lambda: self.status_var.set("🔄 Scanning documents..."))
            
            # Load documents
            texts, file_info = self.read_documents_from_folder(folder)
            if not texts:
                self.root.after(0, lambda: messagebox.showwarning(
                    "Warning", "No supported files found in the selected folder"))
                return
                
            self.supported_files_count = len(file_info)
            self.root.after(0, lambda: self.status_var.set(f"📚 Loaded {len(texts)} documents. Creating embeddings..."))
            self.root.after(0, lambda: self.doc_count_var.set(
                f"📊 Files processed: {file_info['pdf']} PDF, {file_info['docx']} DOCX, {file_info['pptx']} PPTX, {file_info['xml']} XML"))
            
            # Split texts into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )
            
            all_chunks = []
            for text in texts:
                chunks = text_splitter.split_text(text)
                all_chunks.extend(chunks)
                
            self.root.after(0, lambda: self.status_var.set(f"⚙️ Processing {len(all_chunks)} text chunks..."))
            
            # Create QA system
            self.create_qa_system(all_chunks)
            
            # Update UI
            self.root.after(0, lambda: self.progress.stop())
            self.root.after(0, lambda: self.status_var.set(f"🟢 Ready! Loaded {self.supported_files_count} documents successfully."))
            self.root.after(0, lambda: self.query_entry.config(state='normal'))
            self.root.after(0, lambda: self.ask_btn.config(state='normal'))
            self.root.after(0, lambda: self.load_btn.config(state='normal'))
            self.root.after(0, lambda: self.add_to_chat("🤖 System", 
                f"✅ Documents loaded successfully! Processed {self.supported_files_count} files with {len(all_chunks)} text chunks. You can now ask questions!", 
                color='success'))
            
            self.documents_loaded = True
            
        except Exception as e:
            self.root.after(0, lambda: self.progress.stop())
            self.root.after(0, lambda: self.load_btn.config(state='normal'))
            self.root.after(0, lambda: self.status_var.set("❌ Error loading documents"))
            self.root.after(0, lambda: messagebox.showerror("Error", f"Failed to load documents: {str(e)}"))
            
    def read_documents_from_folder(self, folder):
        texts = []
        file_info = {'pdf': 0, 'docx': 0, 'pptx': 0, 'xml': 0}
        supported_files = []
        
        # Scan for supported files
        for file in os.listdir(folder):
            if file.endswith(('.pdf', '.docx', '.pptx', '.xml')):
                supported_files.append(file)
                
        for i, file in enumerate(supported_files):
            full_path = os.path.join(folder, file)
            self.root.after(0, lambda f=file, idx=i+1, total=len(supported_files): 
                          self.status_var.set(f"📖 Reading {f} ({idx}/{total})..."))
            
            try:
                text = ""
                if file.endswith(".pdf"):
                    text = self.read_pdf(full_path)
                    file_info['pdf'] += 1
                elif file.endswith(".docx"):
                    text = self.read_word(full_path)
                    file_info['docx'] += 1
                elif file.endswith(".pptx"):
                    text = self.read_powerpoint(full_path)
                    file_info['pptx'] += 1
                elif file.endswith(".xml"):
                    text = self.read_xml(full_path)
                    file_info['xml'] += 1
                    
                if text.strip():
                    texts.append(f"[Source: {file}]\n{text}")
            except Exception as e:
                print(f"Error reading {file}: {str(e)}")
                continue
                
        return texts, file_info
        
    def read_pdf(self, path):
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
        
    def read_word(self, path):
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    
    def read_powerpoint(self, path):
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text += " | ".join(row_text) + "\n"
                
                text_content.append(slide_text)
            
            return "\n".join(text_content)
        except Exception as e:
            raise Exception(f"Error reading PowerPoint file: {str(e)}")
    
    def read_xml(self, path):
        """Extract text content from XML files"""
        try:
            tree = ET.parse(path)
            root = tree.getroot()
            
            def extract_text_from_element(element, level=0):
                text_parts = []
                indent = "  " * level
                
                # Add element tag and attributes info
                if element.tag and level < 10:  # Prevent too deep nesting
                    tag_info = f"{indent}<{element.tag}"
                    if element.attrib:
                        attrs = ", ".join([f"{k}='{v}'" for k, v in element.attrib.items()])
                        tag_info += f" {attrs}"
                    tag_info += ">"
                    text_parts.append(tag_info)
                
                # Add element text content
                if element.text and element.text.strip():
                    text_parts.append(f"{indent}  {element.text.strip()}")
                
                # Recursively process child elements
                for child in element:
                    child_text = extract_text_from_element(child, level + 1)
                    if child_text:
                        text_parts.append(child_text)
                
                # Add tail text
                if element.tail and element.tail.strip():
                    text_parts.append(f"{indent}{element.tail.strip()}")
                
                return "\n".join(text_parts)
            
            extracted_text = extract_text_from_element(root)
            
            # Also try to extract just text content without structure
            all_text = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    all_text.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    all_text.append(elem.tail.strip())
            
            # Combine structured and unstructured text
            result = f"=== XML Structure ===\n{extracted_text}\n\n=== XML Content ===\n{' '.join(all_text)}"
            return result
            
        except Exception as e:
            raise Exception(f"Error reading XML file: {str(e)}")
        
    def create_qa_system(self, texts):
        if self.embed_model is None:
            self.embed_model = HuggingFaceEmbeddings(model_name='all-MiniLM-L6-v2')
            
        # Create FAISS vector store
        faiss_store = LangchainFAISS.from_texts(texts, self.embed_model)
        
        # Initialize Llama 3.2 model
        llm = Ollama(model="llama3.2")
        
        # Create QA chain
        self.qa_system = RetrievalQA.from_chain_type(
            llm=llm, 
            chain_type="stuff",
            retriever=faiss_store.as_retriever(search_kwargs={"k": 4}),
            return_source_documents=False
        )
        
    def ask_question_event(self, event):
        self.ask_question_thread()
        
    def ask_question_thread(self):
        thread = threading.Thread(target=self.ask_question)
        thread.daemon = True
        thread.start()
        
    def ask_question(self):
        if not self.documents_loaded or not self.qa_system:
            self.root.after(0, lambda: messagebox.showwarning("Warning", "Please load documents first"))
            return
            
        query = self.query_var.get().strip()
        if not query:
            return
            
        try:
            # Record start time
            start_time = time.time()
            
            # Update UI
            self.root.after(0, lambda: self.ask_btn.config(state='disabled'))
            self.root.after(0, lambda: self.status_var.set("🤔 AI is thinking..."))
            self.root.after(0, lambda: self.add_to_chat("👤 You", query, color='user'))
            self.root.after(0, lambda: self.query_var.set(""))
            
            # Get answer
            answer = self.qa_system.run(query)
            
            # Calculate response time
            end_time = time.time()
            response_time = end_time - start_time
            self.response_times.append(response_time)
            
            # Format response time
            if response_time < 1:
                time_str = f"{response_time*1000:.0f}ms"
            else:
                time_str = f"{response_time:.2f}s"
            
            # Update performance metrics
            avg_time = sum(self.response_times) / len(self.response_times)
            if avg_time < 1:
                avg_str = f"{avg_time*1000:.0f}ms"
            else:
                avg_str = f"{avg_time:.2f}s"
            
            perf_text = f"⚡ Last response: {time_str} | Average: {avg_str} | Total queries: {len(self.response_times)}"
            
            # Update UI with answer and performance metrics
            answer_with_time = f"{answer}\n\n⏱️ Response time: {time_str}"
            self.root.after(0, lambda: self.add_to_chat("🤖 AI Assistant", answer_with_time, color='assistant'))
            self.root.after(0, lambda: self.status_var.set("🟢 Ready for next question"))
            self.root.after(0, lambda: self.perf_var.set(perf_text))
            self.root.after(0, lambda: self.ask_btn.config(state='normal'))
            
        except Exception as e:
            self.root.after(0, lambda: self.add_to_chat("❌ System", f"Error: {str(e)}", color='error'))
            self.root.after(0, lambda: self.status_var.set("❌ Error processing question"))
            self.root.after(0, lambda: self.ask_btn.config(state='normal'))
            
    def add_to_chat(self, sender, message, color='default'):
        """Add message to chat with color coding"""
        self.chat_display.config(state='normal')
        
        # Define color schemes
        color_schemes = {
            'user': {'fg': self.colors['accent_blue'], 'bg': None},
            'assistant': {'fg': self.colors['accent_green'], 'bg': None},
            'system': {'fg': self.colors['accent_purple'], 'bg': None},
            'error': {'fg': self.colors['error'], 'bg': None},
            'success': {'fg': self.colors['success'], 'bg': None},
            'default': {'fg': self.colors['text_primary'], 'bg': None}
        }
        
        # Configure text tags for colors
        for tag, colors in color_schemes.items():
            self.chat_display.tag_configure(tag, foreground=colors['fg'])
        
        # Add timestamp
        import datetime
        timestamp = datetime.datetime.now().strftime("%H:%M:%S")
        
        # Insert message with color
        self.chat_display.insert(tk.END, f"\n[{timestamp}] ", 'default')
        self.chat_display.insert(tk.END, f"{sender}:\n", color)
        self.chat_display.insert(tk.END, f"{message}\n", 'default')
        
        self.chat_display.config(state='disabled')
        self.chat_display.see(tk.END)
        
    def clear_chat(self):
        self.chat_display.config(state='normal')
        self.chat_display.delete(1.0, tk.END)
        self.chat_display.config(state='disabled')
        self.add_welcome_message()
    
    def show_performance_stats(self):
        """Show detailed performance statistics"""
        if not self.response_times:
            messagebox.showinfo("Performance Stats", "No queries processed yet!")
            return
        
        stats_window = tk.Toplevel(self.root)
        stats_window.title("📊 Performance Statistics")
        stats_window.geometry("500x400")
        stats_window.configure(bg=self.colors['bg_secondary'])
        stats_window.transient(self.root)
        stats_window.grab_set()
        
        # Center the window
        stats_window.update_idletasks()
        x = (stats_window.winfo_screenwidth() // 2) - (stats_window.winfo_width() // 2)
        y = (stats_window.winfo_screenheight() // 2) - (stats_window.winfo_height() // 2)
        stats_window.geometry(f"+{x}+{y}")
        
        # Create stats content
        stats_frame = ttk.Frame(stats_window, style='Card.TLabelFrame', padding="20")
        stats_frame.pack(fill='both', expand=True, padx=20, pady=20)
        
        # Calculate statistics
        total_queries = len(self.response_times)
        avg_time = sum(self.response_times) / total_queries
        min_time = min(self.response_times)
        max_time = max(self.response_times)
        total_time = sum(self.response_times)
        
        # Format times
        def format_time(t):
            return f"{t*1000:.0f}ms" if t < 1 else f"{t:.2f}s"
        
        # Create stats text
        stats_text = f"""
🚀 AI Performance Statistics

📊 Query Statistics:
   • Total Queries: {total_queries}
   • Average Response Time: {format_time(avg_time)}
   • Fastest Response: {format_time(min_time)}
   • Slowest Response: {format_time(max_time)}
   • Total Processing Time: {format_time(total_time)}

📈 Performance Analysis:
   • Queries under 1 second: {sum(1 for t in self.response_times if t < 1)}
   • Queries 1-3 seconds: {sum(1 for t in self.response_times if 1 <= t < 3)}
   • Queries over 3 seconds: {sum(1 for t in self.response_times if t >= 3)}

📚 Document Information:
   • Documents Loaded: {self.supported_files_count}
   • AI Model: Llama 3.2
   • Embedding Model: all-MiniLM-L6-v2

⚡ System Recommendations:
   {"• Performance is excellent!" if avg_time < 2 else "• Consider optimizing document chunks for faster responses." if avg_time < 5 else "• System may benefit from hardware upgrade or model optimization."}
   {"• Response times are consistent." if max_time - min_time < 2 else "• Response times vary - this is normal for different query complexities."}
        """
        
        # Create scrollable text widget for stats
        text_widget = scrolledtext.ScrolledText(
            stats_frame, wrap=tk.WORD, font=('Consolas', 10),
            bg=self.colors['bg_accent'], fg=self.colors['text_primary'],
            height=20, width=60
        )
        text_widget.pack(fill='both', expand=True)
        text_widget.insert('1.0', stats_text)
        text_widget.configure(state='disabled')
        
        # Close button
        close_btn = ttk.Button(stats_frame, text="✅ Close", 
                              command=stats_window.destroy, style='Success.TButton')
        close_btn.pack(pady=(10, 0))

def main():
    # Check if Ollama is available
    try:
        test_llm = Ollama(model="llama3.2")
        print("✅ Ollama connection successful")
    except Exception as e:
        print(f"⚠️ Warning: Could not connect to Ollama: {e}")
        print("🔧 Setup Instructions:")
        print("   1. Install Ollama: curl -fsSL https://ollama.ai/install.sh | sh")
        print("   2. Pull model: ollama pull llama3.2")
        print("   3. Start Ollama service")
        
    # Create and run the application
    root = tk.Tk()
    
    # Set application icon (if available)
    try:
        root.iconbitmap('chatbot.ico')  # Add your own icon file
    except:
        pass
    
    app = ModernDocumentChatGUI(root)
    
    # Center window on screen after app initialization to get correct dimensions
    # root.update_idletasks()
    # screen_width = root.winfo_screenwidth()
    # screen_height = root.winfo_screenheight()
    # window_width = root.winfo_width()
    # window_height = root.winfo_height()
    
    # x = (screen_width // 2) - (window_width // 2)
    # y = (screen_height // 2) - (window_height // 2)
    
    # # Ensure window doesn't go off-screen
    # x = max(0, min(x, screen_width - window_width))
    # y = max(0, min(y, screen_height - window_height))
    
    # root.geometry(f"+{x}+{y}")
    
    # root.mainloop()
    
    def center_window(root, width=800, height=600):
        # Set the window size first
        root.geometry(f"{width}x{height}")
        
        # Update to get accurate dimensions
        root.update_idletasks()
        
        # Get screen dimensions
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        
        # Calculate center position
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        
        # Ensure window doesn't go off-screen
        x = max(0, min(x, screen_width - width))
        y = max(0, min(y, screen_height - height))
        
        # Set final position
        root.geometry(f"{width}x{height}+{x}+{y}")
        
    center_window(root, 800, 600)
    root.mainloop()

if __name__ == "__main__":
    main()